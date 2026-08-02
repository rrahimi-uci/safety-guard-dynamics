#!/usr/bin/env python
"""Build the executive deck: should we buy a hosted guardrail, or run our own?

    python slides/make_exec_deck.py       (from papers/unified-report/)
    -> slides/safety_guard_exec_deck.pptx  (16:9, ~10 slides, speaker notes)

This is NOT a shortened research deck. `make_deck.py` answers "what did we measure and is
it sound"; this answers "what should we do, what will it cost, and what are we still unsure
about". Different question, so different structure: the recommendation is on slide 3 rather
than slide 18, method appears only where it changes whether you believe a number, and every
confidence interval lives in the speaker notes instead of the slide body.

Numbers come from `frontier_numbers.load()`, which parses the same
`generated/frontier_macros.tex` the report `\\input`s, so a slide figure cannot drift from
the paper. Percentages are rounded for readability; the underlying decimals and intervals
are in the notes so a challenge from the room can be answered exactly.

Visual register is deliberately larger and emptier than the research deck: an executive
audience reads a slide in about four seconds, so each one carries a single claim, and the
claim is in the headline rather than in the body.
"""
# ruff: noqa: E741
#   `l` as a left-coordinate parameter is the convention make_deck.py already uses across
#   every geometry helper. Matching it keeps the two deck scripts readable side by side,
#   which matters more here than the ambiguous-name rule.
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
sys.path.insert(0, str(HERE))
import frontier_numbers as FN  # noqa: E402

# Act I re-read over FPR [0,.05], from the macro file the paper \inputs -- so the exec deck
# cannot quote a different amplification than Table 6 of the report.
LFP = FN.load_named("lowfpr_macros.tex")

OUT = HERE / "safety_guard_exec_deck.pptx"

# ─────────────────────────────────────────────────────────── identity (shared palette)
# Tokens come from deck_theme, the same module the benchmark deck and both figure generators
# use. The exec redesign was verified to use an identical palette and type scale; only its
# geometry differs (a wider 0.85" margin and a 40pt hero), so those stay local below.
import deck_theme as T  # noqa: E402

INK = T.rgb(T.TEXT)
SLATE = T.rgb(T.DIM)
MUTED = T.rgb(T.BODY)
RULE = T.rgb(T.CARD_LINE)
PAPER = T.rgb(T.CARD)
WHITE = T.rgb(T.TEXT)
BG_TITLE = T.rgb(T.BG_TITLE)
BG_SLIDE = T.rgb(T.BG_SLIDE)
CARD_LINE = T.rgb(T.CARD_LINE)
WARN_CARD = T.rgb(T.WARN_CARD)
WARN_LINE = T.rgb(T.WARN_LINE)
FAINT = T.rgb(T.FAINT)
DATA = T.rgb(T.DATA)
ACCENT = T.rgb(T.ACCENT)
ACCENT_SOFT = T.rgb(T.ACCENT_SOFT)
BLUE = T.rgb(T.DATA_REPRESENTED)
GREEN = T.rgb(T.DATA_COMPOSITION)
AMBER = T.rgb(T.DATA_GOLD)
# On a dark surface a tint is a surface, not a wash: red-keyed panels take the warning
# surface, everything else the standard card.
TINT = {ACCENT: WARN_CARD, ACCENT_SOFT: WARN_CARD, BLUE: PAPER, GREEN: PAPER,
        AMBER: PAPER, SLATE: PAPER}
TINT_LINE = {ACCENT: WARN_LINE, ACCENT_SOFT: WARN_LINE, BLUE: CARD_LINE, GREEN: CARD_LINE,
             AMBER: CARD_LINE, SLATE: CARD_LINE}
LABEL_ON_CARD = {ACCENT: ACCENT_SOFT}   # red is too low-contrast as a label on the dark card
SERIF, SANS = T.DISPLAY, T.UI
W, H = Inches(T.SLIDE_W), Inches(T.SLIDE_H)
M = Inches(0.85)                        # exec geometry: wider margin than the benchmark deck
CW = W - 2 * M                          # 11.633
TITLE_SHORT = "Guardrail sourcing · in-house inline, escalate the uncertain slice"


# ───────────────────────────────────────────────────────────────────────── primitives
def _noshadow(shape):
    spPr = shape.fill._xPr
    for tag in ("a:effectLst", "a:effectDag"):
        for el in spPr.findall(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}"):
            spPr.remove(el)


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, int(l), int(t), int(w), int(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    _noshadow(sh)
    sh.text_frame.word_wrap = True
    return sh


def tbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, first=False, align=PP_ALIGN.LEFT, space_after=8, line_spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    return p


def run(p, text, size=16, bold=False, color=INK, font=SANS, italic=False, spc=None):
    # One run per line with an <a:br/> between, for the reason make_deck.run documents: a
    # literal "\n" inside <a:t> is not an OOXML line break and renderers may drop it.
    r = None
    for i, segment in enumerate(str(text).split("\n")):
        if i:
            p.add_line_break()
        r = p.add_run()
        r.text = segment
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
        if spc is not None:
            # Letter-spacing has no python-pptx property; set the raw attribute, exactly as
            # make_deck._spacing does. Note qn() is wrong here -- it expects a prefixed tag.
            r.font._rPr.set("spc", str(int(spc)))
    return r


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.n = 0

    def _background(self, s, fill):
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = fill

    def blank(self, chrome=True, title_slide=False):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._background(s, BG_TITLE if title_slide else BG_SLIDE)
        # No top accent bar: the redesign removed it (deck_theme.HAS_TOP_BAR).
        if chrome:
            self.n += 1
            tf = tbox(s, M, Inches(T.FOOTER_Y), Inches(9.5), Inches(0.28))
            p = para(tf, first=True, space_after=0)
            run(p, TITLE_SHORT, size=T.SZ_FOOTER, color=FAINT)
            tf2 = tbox(s, W - M - Inches(1.2), Inches(T.FOOTER_Y), Inches(1.2), Inches(0.28))
            p2 = para(tf2, first=True, align=PP_ALIGN.RIGHT, space_after=0)
            run(p2, f"{self.n:02d}", size=T.SZ_PAGENUM, bold=True, color=ACCENT, spc=60)
        return s

    def header(self, s, kicker, title, sub=None):
        tf = tbox(s, M, Inches(0.40), CW, Inches(0.28))
        p = para(tf, first=True, space_after=0)
        run(p, kicker.upper(), size=10.5, bold=True, color=ACCENT, spc=140)
        # 1.02" of title box holds two 25pt lines, which several exec titles need.
        tf = tbox(s, M, Inches(0.72), CW, Inches(1.02))
        p = para(tf, first=True, space_after=0, line_spacing=1.02)
        run(p, title, size=T.SZ_TITLE, bold=True, color=INK, font=SERIF)
        y = Inches(2.10)
        if sub:
            tf = tbox(s, M, Inches(1.86), CW, Inches(0.34))
            p = para(tf, first=True, space_after=0)
            run(p, sub, size=12.5, color=SLATE)
            y = Inches(2.45)
        # No rule under the header: the redesign dropped it.
        return y

    def notes(self, s, text):
        s.notes_slide.notes_text_frame.text = text.strip()

    def save(self):
        # See make_deck.Deck.save: replace python-pptx's default-template document
        # properties, and correct the inherited 4:3 PresentationFormat on a 16:9 deck.
        T.stamp_properties(
            self.prs, "Guardrail sourcing: buy hosted, or run our own?",
            subject="Executive briefing accompanying “Benchmark Gains Do Not Guarantee "
                    "Safety Transfer: A Comprehensive Study of Fine-Tuning Small Language Model Safety Guards for High-Compliance and General Safety Domains”")
        self.prs.save(OUT)
        T.fix_presentation_format(OUT)
        return OUT


def statcard(s, l, t, w, h, value, caption, color=ACCENT, value_size=None):
    """Flat card with a display numeral. No top accent rule: the border carries the edge."""
    value_size = T.SZ_STAT if value_size is None else value_size
    rect(s, l, t, w, h, fill=TINT[color], line=TINT_LINE[color], lw=T.LW_CARD)
    tf = tbox(s, l + Inches(0.26), t + Inches(0.16), w - Inches(0.52), h - Inches(0.32))
    p = para(tf, first=True, space_after=4)
    run(p, value, size=value_size, bold=True, color=color, font=SERIF)
    p = para(tf, space_after=0)
    run(p, caption, size=T.SZ_BODY, color=MUTED)


def bullets(s, l, t, w, h, items, size=None, gap=12, mcolor=ACCENT):
    size = 11.0 if size is None else size      # exec body size in the redesign
    tf = tbox(s, l, t, w, h)
    for i, item in enumerate(items):
        p = para(tf, first=(i == 0), space_after=gap)
        run(p, "—  ", size=size, bold=True, color=mcolor)
        if isinstance(item, tuple):
            run(p, item[0], size=size, bold=True, color=INK)
            run(p, item[1], size=size, color=MUTED)
        else:
            run(p, item, size=size, color=MUTED)


def table(s, l, t, w, rows, col_w, row_h=Inches(0.46), header=True):
    """Minimal table: first row is a header band, subsequent rows alternate paper/white."""
    y = t
    for ri, row in enumerate(rows):
        fill = PAPER if (header and ri == 0) else (BG_SLIDE if ri % 2 else PAPER)
        rect(s, l, y, w, row_h, fill=fill, line=CARD_LINE, lw=T.LW_CARD)
        x = l
        for ci, cell in enumerate(row):
            cwid = int(w * col_w[ci])
            tf = tbox(s, x + Inches(0.14), y, cwid - Inches(0.28), row_h,
                      anchor=MSO_ANCHOR.MIDDLE)
            p = para(tf, first=True, space_after=0,
                     align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            bold = (header and ri == 0) or ci == 0
            col = DATA if (header and ri == 0) else MUTED
            txt, cc = (cell if isinstance(cell, tuple) else (cell, col))
            run(p, str(txt), size=T.SZ_KICKER_SM, bold=bold, color=cc)
            x += cwid
        y += row_h
    return y


def picture(s, name, l, t, w, h, align="center"):
    """Place a figure inside an (l,t,w,h) box, preserving aspect ratio.

    The figures are generated at their own aspect by make_exec_figures.py, so fitting
    them to a fixed box would stretch them. Scale to fit, then align within the box.
    """
    from PIL import Image
    path = ASSETS / f"{name}.png"
    if not path.is_file():
        raise FileNotFoundError(
            f"{path} missing -- run `python slides/make_exec_figures.py` first")
    iw, ih = Image.open(path).size
    scale = min(w / iw, h / ih)
    dw, dh = int(iw * scale), int(ih * scale)
    if align == "left":
        x = l
    elif align == "right":
        x = l + w - dw
    else:
        x = l + (w - dw) // 2
    y = t + (h - dh) // 2
    return s.shapes.add_picture(str(path), int(x), int(y), dw, dh)


def callout(s, l, t, w, h, label, body, color=ACCENT):
    rect(s, l, t, w, h, fill=TINT[color], line=TINT_LINE[color], lw=T.LW_CARD)
    tf = tbox(s, l + Inches(0.26), t + Inches(0.16), w - Inches(0.52), h - Inches(0.32))
    p = para(tf, first=True, space_after=5)
    run(p, label.upper(), size=T.SZ_BODY_SM, bold=True,
        color=LABEL_ON_CARD.get(color, color), spc=110)
    p = para(tf, space_after=0, line_spacing=1.06)
    run(p, body, size=11.0, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════ build
F = FN.load()
HT = FN.load_h2h()
CA = FN.load_cascade()
b = FN.bare
d = Deck()

# ─────────────────────────────────────────────────────────────── 1 · title
s = d.blank(chrome=False, title_slide=True)

# Concentric rings, right of the headline — the redesign's title motif, sized for the exec
# deck's wider margin. Drawn first so the text layers sit above them.
for cx, cy, dia, ln, lw in [(7.10, 1.30, 4.90, T.DECO_LINE, T.LW_DECO),
                            (7.75, 1.95, 3.60, T.DECO_LINE, T.LW_DECO),
                            (8.40, 2.60, 2.30, T.DECO_LINE_WARM, T.LW_ACCENT_THIN)]:
    rect(s, Inches(cx), Inches(cy), Inches(dia), Inches(dia), fill=None,
         line=T.rgb(ln), lw=lw, shape=MSO_SHAPE.OVAL)
rect(s, Inches(9.41), Inches(3.61), Inches(0.28), Inches(0.28), fill=ACCENT,
     shape=MSO_SHAPE.OVAL)

tf = tbox(s, M, Inches(2.12), Inches(8.20), Inches(0.30))
p = para(tf, first=True, space_after=0)
run(p, "GUARDRAIL SOURCING DECISION", size=11.0, bold=True, color=ACCENT, spc=180)
tf = tbox(s, M, Inches(2.56), Inches(8.60), Inches(1.90))
p = para(tf, first=True, space_after=0, line_spacing=1.03)
run(p, "Do we buy a hosted guardrail,\nor run our own?", size=40, bold=True,
    color=INK, font=SERIF)
tf = tbox(s, M, Inches(4.62), Inches(7.70), Inches(0.80))
p = para(tf, first=True, space_after=0, line_spacing=1.10)
run(p, "Measured on 2,275 expert-annotated finance, healthcare and law prompts. "
       "Twenty-four guard configurations, identical rows, identical false-alarm budget.",
    size=T.SZ_SUBTITLE, color=MUTED)
tf = tbox(s, M, Inches(6.02), Inches(7.70), Inches(0.30))
p = para(tf, first=True, space_after=0)
run(p, "Reza Rahimi, PhD", size=12.5, bold=True, color=INK)
tf = tbox(s, M, Inches(6.32), Inches(7.70), Inches(0.28))
p = para(tf, first=True, space_after=0)
run(p, "Full method, intervals and limitations: “Safety Benchmark Gains Do Not Guarantee Safety Transfer: "
       "A Comprehensive Study of Fine-Tuning Small Language Model Safety Guards for High-Compliance and General Safety Domains” — the technical report and research deck",
    size=10.5, color=SLATE)
d.notes(s, """
Ninety seconds on this slide. The framing sentence is: we need a safety guardrail in front of
every request our assistant handles, and there are two ways to get one -- call a hosted
frontier model, or run a small open-weights model ourselves. This deck prices that choice.

Say up front what makes the comparison trustworthy, because it is the thing that makes the rest
land: every configuration is scored on the SAME rows, and every one is re-tuned to the SAME
false-alarm rate before we compare catch rates. Without the second step you can make any guard
look good by letting it cry wolf more often.

Do not promise a single recommendation yet. Slide 3 gives it.
""")

# ────────────────────────────────────────────────────── 2 · why this is a decision
s = d.blank()
y = d.header(s, "the decision", "A guardrail runs on every request, so its cost is a"
             " recurring line, not a one-off",
             "Which means accuracy, latency, unit cost and data residency all bind at once")
cw = (CW - Inches(0.4)) / 2
bullets(s, M, y + Inches(0.06), cw, Inches(3.4), [
    ("Accuracy. ", "A missed unsafe prompt is the incident. A false alarm is a blocked "
     "customer. Both are real costs, so the guard must be judged at a fixed alarm rate."),
    ("Latency. ", "It sits in the request path. Whatever it adds, every user pays on every "
     "turn."),
    ("Unit cost. ", "Priced per prompt, it scales with traffic, not with headcount."),
    ("Residency. ", "In mortgage the prompt itself contains borrower information, so where "
     "it travels is a compliance question before it is an engineering one."),
], size=14.5)
callout(s, M + cw + Inches(0.4), y + Inches(0.06), cw, Inches(1.80),
        "why we could not just read the vendor benchmarks",
        "Public guard leaderboards score general web-safety prompts. Our exposure is "
        "regulated-domain advice, where the same words are safe or unsafe depending on who "
        "is asking and what was promised.", color=BLUE)
# NOT "on our own terms". ExpGuard is a third-party, expert-annotated set -- which is the
# strength of the evidence, and the phrase invited a legal or risk stakeholder to hear
# "we evaluated on company traffic", which is the opposite of what happened.
callout(s, M + cw + Inches(0.4), y + Inches(2.04), cw, Inches(1.80),
        "so we measured on external, expert-annotated prompts",
        "Third-party expert-annotated finance, healthcare and law prompts — not our own "
        "traffic and not our own labels; every guard scored on identical rows; catch rates "
        "compared only at a matched 5% false-alarm budget.",
        color=GREEN)
d.notes(s, """
Purpose of this slide: establish that this is a sourcing decision with four binding
constraints, not a model-quality question. If someone in the room only cares about accuracy,
this is where you widen it.

The residency bullet is the one that tends to change the conversation with legal. Our mortgage
prompts contain borrower detail; sending them to a third party is a GLBA question. That is not
a reason to refuse hosted models, but it is a reason the decision cannot be made by the ML team
alone.

If asked why not just trust published guard benchmarks: because they measure general web
safety, and our failure mode is regulated-domain advice where context decides the label. We
have direct evidence in the technical report that guard rankings reorder when the benchmark
changes. Name that correctly if you quote it: the report's TITLE finding is that a gain on a
represented benchmark does not establish transfer. The ranking reorder across benchmarks is
the corollary, and it is the one that bears on a sourcing decision.

Be exact about "our own terms" if anyone in legal or risk pushes on the second callout. The
evaluation ran on a third-party, expert-annotated set. No company traffic, no borrower data,
and no labels we produced ourselves left this building or entered the benchmark.
""")

# ──────────────────────────────────────────────────────── 3 · the recommendation
s = d.blank()
y = d.header(s, "the answer", "Run the small guard in-house on every request. Send out only "
             "the ones it is unsure about",
             "One rule, three lanes. The question is not in-house OR outsource — it is which "
             "requests go out")
rows = [["Lane", "Which requests", "Guard", "Why"],
        [("1  Regulated", INK), "prompt contains borrower or patient data",
         ("in-house only", GREEN), "cannot lawfully leave"],
        [("2  Uncertain", INK), "in-house guard is near its decision line",
         ("escalate to hosted", ACCENT), "this is where the accuracy is won"],
        [("3  Clear-cut", INK), "in-house guard is confident either way",
         ("in-house only", GREEN), "hosted would not change the answer"]]
table(s, M, y + Inches(0.04), CW, rows, [0.17, 0.34, 0.21, 0.28], row_h=Inches(0.62))
cw3 = (CW - Inches(0.6)) / 3
statcard(s, M, y + Inches(2.72), cw3, Inches(1.50), "20%",
         "of requests escalated — lane 2 only", color=ACCENT)
statcard(s, M + cw3 + Inches(0.3), y + Inches(2.72), cw3, Inches(1.50), "51%",
         "of the accuracy gap to hosted, closed", color=GREEN)
statcard(s, M + 2 * (cw3 + Inches(0.3)), y + Inches(2.72), cw3, Inches(1.50), "80%",
         "of prompts never leave our network", color=BLUE)
# ONE notes call. There were two, and `notes()` assigns rather than appends, so the second
# silently discarded the first -- the deck shipped without the lane walkthrough, and with a
# notes block whose numbers derive none of the three figures on the slide.
d.notes(s, f"""
This is the slide. If they remember one thing, it is that the question was posed wrongly: it
is not in-house versus outsource, it is which requests go out.

Walk the three lanes in order. Lane 1 is a legal constraint, not a technical choice -- if the
prompt carries borrower or patient data it stays in, full stop. Lane 3 is the cheap majority:
when the in-house guard is confident, a second opinion changes nothing, so paying for one is
waste. Lane 2 is the whole game -- the band around the in-house guard's decision line, where it
is genuinely unsure, and where the hosted model's advantage actually lives.

WHERE THE THREE NUMBERS COME FROM, because they are one measurement and not three, and because
the obvious interval to quote is the wrong one:
  - 51% = ({b(CA['Twenty'])} - {b(CA['LocalOnly'])}) / ({F['BestTpr']} - {b(CA['LocalOnly'])}),
    i.e. the SmolLM3-3B inline guard against {F['BestName']} at a matched 5% budget, with the
    least-confident fifth escalated. The gap being closed is {F['GainOverBase']} -- that pair.
    Computed on the underlying scores that is 51.1%; do the same sum with the three-decimal
    published values and you get 50.5%. Rounding, not a discrepancy -- same for the 64% on
    slide 8.
  - The {F['GainOverOpen']} {F['GainOverOpenCI']} figure is a DIFFERENT pair (hosted against
    {F['BestOpenName']}, our strongest open guard) and belongs on the next slide. Do not quote
    it as the interval behind 51%.
  - 20% and 80% are the same dial read two ways: escalate one request in five, four in five
    never leave the network.

Slide 8 has the curve, and it is smooth -- you can dial the escalated share to whatever your
legal and budget constraints allow rather than choosing an architecture.

The other figures, precisely, for anyone who asks:
  - {F['Slowdown']}x is median-to-median: about {F['BestMedianMs']} ms hosted against
    10-25 ms for a small guard on our own GPU.
  - ${F['BestCost']} per thousand prompts is billed tokens at public list prices, and it
    excludes the GPU you already own on the self-hosted side. It is an estimate, not an
    invoice. We have not amortised our own side; say so rather than improvising a number.

If asked "why not escalate everything": you would get the full {FN.pct(F['BestTpr'])}, and you
would also send every borrower prompt to a third party, pay ${F['BestCost']} per thousand, and
put {F['BestMedianMs']} ms on every request. Lane 1 forbids it anyway for the traffic we care
most about.

If asked "why not escalate nothing": that is the {FN.pct(b(CA['LocalOnly']))} floor for the
inline 3B guard, and it means knowingly missing about one unsafe prompt in ten that we could
have caught.

The recommendation is deliberately a split, not a winner. If pushed for one answer: hosted,
because the accuracy difference is the largest single effect anywhere in this study. But the
regulated path is exactly where we cannot take it, which is why the split is the honest answer
rather than a hedge.

Someone will propose "just fine-tune ours" or "just use a bigger one". Both were tested.
Neither worked. Slides 6 and 7.
""")

# ─────────────────────────────────────────────────────── 4 · what hosted buys
s = d.blank()
y = d.header(s, "what hosted buys", "At the same false-alarm rate, the hosted model catches "
             "about half of the unsafe prompts our own guards miss",
             "Catch rate on expert-annotated finance / healthcare / law prompts, "
             "all at a 5% false-alarm budget")
picture(s, "exec_gap", M, y - Inches(0.06), Inches(8.05), Inches(4.30), align="left")
# "The true gap is if anything wider" was withdrawn from the report. Coarse ties bound how
# finely the hosted ranking resolves; they do not fix a direction, so the gap is not a floor.
callout(s, M + Inches(8.35), y + Inches(0.02), CW - Inches(8.35), Inches(1.86),
        "read this as coarsely resolved",
        "The hosted model reports only a coarse 0-100 confidence, which limits how finely we "
        "can rank its answers. That bounds precision without fixing a direction — a finer "
        "score could order a tied block either way — so do not read the gap as conservative.",
        color=BLUE)
callout(s, M + Inches(8.35), y + Inches(2.06), CW - Inches(8.35), Inches(1.86),
        "why 'at the same false-alarm rate' matters",
        "Any guard can raise its catch rate by alarming more often. Holding the alarm rate "
        "fixed at 5% is what makes every bar here comparable at all.", color=AMBER)
d.notes(s, f"""
One claim: hosted is better, by about ten points of catch rate, and the comparison is fair.

Exact figures if challenged: {F['BestName']} {F['BestTpr']} against {F['BestOpenName']}
{F['BestOpenTpr']}, difference {F['GainOverOpen']} with 95% interval {F['GainOverOpenCI']}.
Against the strongest model in our locked research panel it is {F['GainOverBase']}
{F['GainOverBaseCI']}. Both intervals exclude zero. These are paired comparisons on identical
rows, which removes row-difficulty as an explanation.

The coarse-confidence point is worth making unprompted, because it is the one place we
disadvantaged the hosted model and it strengthens rather than weakens the conclusion.

If asked about the labels: external, expert-annotated, and a stronger labelling tier than
anything we produced ourselves. We did not grade our own homework here.
""")

# ────────────────────────── 4b · the gap is a regime, not a verdict on model size
# The single most decision-relevant slide in this deck: it converts "hosted is better" into
# "hosted is better on traffic we cannot anticipate", which is a sourcing rule rather than a
# preference. Numbers come from h2h_macros.tex via frontier_numbers.load_h2h(), the same
# anti-drift contract the rest of the frontier material uses.
s = d.blank()
# "our own guard" -> "our own TUNED guards". The qualifier is load-bearing, not stylistic:
# include the base arms in the same aggregate and the advantage reverses to -0.264, excluding
# zero in the opposite direction. The body text always said "tuned"; the headline did not.
y = d.header(s, "but better at what",
             "On traffic we can describe in advance, our own tuned guards already win",
             "Same five corpora, same " + b(HT['Budget']) + " false-alarm budget, "
             "split by whether the source was in our training manifest")

cw = (CW - Inches(0.40)) / 2
statcard(s, M, y + Inches(0.10), cw, Inches(1.62),
         b(HT['AggDeltaTpr']) + " catch rate",
         # The interval is conditional on the three corpora we tested; the report states the
         # unconditional version in its abstract, and it includes zero. An exec reading this
         # as "we win on describable traffic" needs the scope on the same card as the claim.
         "Averaged across every source we can describe in advance, our tuned guards rank "
         "BETTER than the hosted model at the same alarm budget. 95% interval "
         + b(HT['AggDeltaTprCI']) + " -- excludes zero across these three sources; treat "
         "them as a sample of future manifests and it widens to "
         + b(HT['AggSrcResampledCI']) + ".",
         color=BLUE)
statcard(s, M + cw + Inches(0.40), y + Inches(0.10), cw, Inches(1.62),
         b(HT['TransRefTpr']) + "  vs  " + b(HT['TransBestLocalTpr']),
         "The same comparison on a source we did NOT train on. Hosted leads, and our best "
         "guard there is an UNTUNED base -- tuning is what costs us the position.",
         color=ACCENT)

# This callout used to recommend escalating on UNFAMILIARITY "not on how unsure the local guard
# sounds". That is a router we never built: the only cascade measured anywhere in this programme
# escalates by rank distance to the local guard's own decision line -- a margin router -- and it is
# what slide 3 recommends and slide 8 prices. Recommending the untested alternative here
# contradicted both the deck and the report, so the callout now recommends what was measured and
# names the alternative as untested.
callout(s, M, y + Inches(1.86), CW, Inches(1.14), "the sourcing rule this implies",
        "The gap is not a ceiling a bigger local model would break through -- it is the price "
        "of the regime. So the question is not “can a small guard match the frontier” but "
        "“what share of our traffic can we describe in a training manifest”. Self-host that "
        "share and escalate the rest by how near the local guard sits to its own decision line "
        "-- the router we measured. Routing on UNFAMILIARITY instead is the more attractive idea "
        "and we did not test it: treat it as the next experiment, not as the plan.",
        color=GREEN)
# The weighting that reverses the sign belongs ON the slide, not only in the notes. It is what
# converts "run a small guard in-house on describable traffic" into "run a small guard we have
# TUNED on a manifest of that traffic" -- a funded project with a maintenance burden, which the
# room has to price before it approves lane 1.
callout(s, M, y + Inches(3.10), CW, Inches(0.94), "what the advantage is a property of",
        "Include the untuned base arms in the same average and it reverses to "
        + b(HT['AggWithBaseTpr']) + " " + b(HT['AggWithBaseTprCI']) + " — excluding zero in "
        "the OPPOSITE direction. So this is not “small guards beat hosted models”; it is "
        "“guards we have tuned on a manifest of that traffic do”. Lane 1 is therefore a "
        "standing tuning programme, not a default configuration.", color=ACCENT)

d.notes(s, f"""
This is the slide that changes the recommendation, so do not rush it.

Slide 4 said hosted is about ten points better. True -- on prompts nobody anticipated. This
slide says that on sources we can enumerate in a training manifest the ordering inverts, and a
our tuned guards rank BETTER on average: the equal-source mean difference is
{b(HT['AggDeltaTpr'])}, 95% interval {b(HT['AggDeltaTprCI'])}, which excludes zero. In AP it is
{b(HT['AggDeltaAp'])} {b(HT['AggDeltaApCI'])}.

Be careful with single benchmarks here, and say this before anyone asks. The most eye-catching
cell is a {b(HT['BestName'])} guard at {b(HT['BestTpr'])} against {b(HT['RefTpr'])} on
{b(HT['BestSource'])}. That cell is the largest of {b(HT['NCells'])}, so quoting its own interval
overstates it: {b(HT['NSigNominal'])} of {b(HT['NCells'])} cells clear zero on their own, and
{b(HT['NSigHolm'])} survive a familywise correction. The average is the number to defend, not
the best cell.

And be precise about what the average is an average OVER, because it bounds the sourcing rule.
The interval is conditional on these three corpora: it resamples rows within a fixed source set,
because three sources we picked do not sample a population of sources. Drawing sources with
replacement instead widens it to {b(HT['AggSrcResampledCI'])}, which INCLUDES zero. Read that as
the honest scope of the claim -- it says our tuned guards rank better on traffic that looks like
these three manifests, not that they would on any manifest we might write next. It is a reason
to characterise our own traffic before assuming the advantage transfers to it, not a reason to
discount the split recommendation, which also rests on residency and cost.

The report tabulates four defensible weightings of these twelve cells and only two of them
support a positive advantage at all. Weighting sources by row count roughly halves it, to
{b(HT['AggRowWTpr'])} {b(HT['AggRowWTprCI'])}, which straddles zero. And the one on the red
callout is the one to keep in view: include the BASE arms alongside the tuned ones and the
aggregate becomes {b(HT['AggWithBaseTpr'])} {b(HT['AggWithBaseTprCI'])}, excluding zero in the
opposite direction. That is not a statistical footnote, it is a budget consequence. The
represented-source advantage is a property of TUNED guards specifically -- the specialization
finding from slide 6 seen from the other side -- so lane 1 is "run a small guard we have tuned
on a manifest of our own traffic, and keep tuning it", which is a funded workstream with a
maintenance burden. Price that before approving lane 1.

Represented sources here are {b(HT['RepSources'])}. Held-out sources are the ones we never
trained on, and there the hosted model leads.

Two honest caveats, offered before they are asked. Per-source samples are small (67 to 451
rows), so read the direction rather than the exact margin. And this is retrospective analysis
on a panel we inspected while building the method -- it is a strong signal, not a sealed
confirmatory result.

The business consequence: we do not need to beat the frontier everywhere to stop paying per
prompt. We need to know which slice of traffic we can characterise, self-host that, and buy
hosted capacity only for the rest.
""")

# ────────────────────────────────────────────────── 5 · tuning did not close it
s = d.blank()
# The subtitle used to say "change in CATCH RATE" over a chart whose axis says AP points and
# whose values are macro-AP on our own internal panel, while the notes narrated a third
# quantity (ExpGuard matched-budget recall) and called the chart "external". Three
# measurements, one label. The report is emphatic that the three flavors are never pooled, so
# the slide now names the one it plots and keeps the external result in its own labelled place.
y = d.header(s, "route one: fine-tune ours", "Fine-tuning did not close the gap, and on "
             "unfamiliar traffic it usually made things worse",
             f"Change in ranking accuracy after fine-tuning (macro-AP), on our own internal "
             f"{F['SftNumTotal']}-checkpoint panel · {F['NSeeds']} training runs each")
picture(s, "exec_tax", M, y - Inches(0.06), Inches(8.15), Inches(4.20), align="left")
callout(s, M + Inches(8.45), y + Inches(0.02), CW - Inches(8.45), Inches(2.18),
        "the pattern behind it",
        "Fine-tuning teaches a model the traffic you trained on and costs it accuracy on "
        "traffic you did not. Read downward: the blue gain shrinks steadily as the base gets "
        "stronger. The red cost is less orderly — it turns from a gain into a loss and stays "
        "there, because every tuned model lands on much the same finishing accuracy.",
        color=ACCENT)
# The external result gets its own labelled home here rather than being narrated over the
# chart. It is a different instrument on different rows -- catch rate at a matched budget on
# expert-annotated prompts, not macro-AP on our panel -- and the report never pools the two.
callout(s, M + Inches(8.45), y + Inches(2.34), CW - Inches(8.45), Inches(2.12),
        "different measurement, same verdict",
        f"On the external expert-annotated prompts — catch rate at a matched 5% budget, never "
        f"pooled with the chart — {F['SftNumHurt']} of {F['SftNumTotal']} checkpoints got "
        f"worse. So budgeting a tuning project to close the vendor gap is not supported: the "
        f"best tuned result anywhere ({FN.pct(F['BestSftTpr'])}) still loses to hosted "
        f"({FN.pct(F['BestTpr'])}).", color=AMBER)
d.notes(s, f"""
KEEP THE TWO MEASUREMENTS APART, because this slide carries both and the report never pools
them. The CHART is macro-AP -- average ranking quality over the whole score list -- measured on
our own internal {F['SftNumTotal']}-checkpoint panel, the same panel the research was run on.
The blue callout is a different instrument: catch rate at a matched 5% false-alarm budget on
the EXTERNAL expert-annotated prompts. Both say fine-tuning does not close the gap. Neither is
a restatement of the other, and averaging them would be a mistake.

On the external measurement the headline is the sign split, not the average. If you quote only
the mean change ({F['SftMeanDelta']}) you will mislead the room -- it is a small number sitting
on top of swings from {F['SftWorstDelta']} to {F['SftBestDelta']}, and
{F['SftNumHurt']} of {F['SftNumTotal']} checkpoints moved the wrong way.

One more number for this slide if the room asks how much worse it is where we would actually
run: read only inside the 5% false-alarm budget, the held-out loss on the internal panel is
about {FN.bare(LFP['LowFprTransAmplification'])}x bigger than the average metric shows. That is
the internal panel again, not the external set.

Mechanism, in plain terms: fine-tuning on a fixed set of sources specialises the model to
those sources. It buys a lot of accuracy on traffic that looks like the training data and
gives some back on traffic that does not. The strongest model we tuned gained the least --
that part is orderly, all the way down the chart.

Do not go one step further and say it also lost the most, which an earlier version of this
slide did. The loss is not ordered by starting strength: the 4B model gave back the most
(-0.150) from a weaker starting point than either the 8B (-0.101) or the 32B (-0.117). What
is orderly is where they all END UP -- every tuned model lands on roughly the same
unfamiliar-traffic accuracy whatever it started from, so the cost is the distance it had to
fall, not evidence that better models specialise harder. The planning consequence is the same
either way: above a strong starting point, tuning buys very little and still charges close to
full price.

This is the central finding of the underlying research. The chart is that finding on the panel
it was measured on; the blue callout is the same pattern recurring on external expert-annotated
data, which is the harder test. Say which one you are quoting -- an earlier version of this
slide described the chart itself as the external result, and it is not.

If asked "did you tune it badly": five independent training runs per model, one recipe, and
the same recipe that produces large gains on in-distribution traffic. The gains are real; they
just do not transfer.
""")

# ─────────────────────────────────────────────── 6 · scale did not close it either
s = d.blank()
y = d.header(s, "route two: buy a bigger one", "Eight times the parameters bought less than "
             "the gap that was left",
             "Catch rate against model size, same family, same prompt, same rows")
picture(s, "exec_scale", M, y - Inches(0.10), Inches(8.10), Inches(3.05), align="left")
bullets(s, M + Inches(8.40), y + Inches(0.04), CW - Inches(8.40), Inches(3.0), [
    (f"{F['ScaleFactor']}× the parameters bought "
     f"{FN.points(F['ScaleGain'])}. ", "Going from 4B to 32B."),
    (f"The gap left is {FN.points(F['GainOverOpen'])}. ",
     "Slightly larger than everything that scaling bought."),
    ("4B to 8B bought nothing. ", "Size does not buy guard accuracy smoothly."),
], size=13.5, gap=11)
# The callout used to read "would take at least another order of magnitude". The report
# withdrew that sentence: three points, one of them non-monotonic, identify no scaling law,
# so no required parameter count can be extrapolated from this ladder. Say what was measured.
callout(s, M, y + Inches(3.10), CW, Inches(0.98), "what this implies for a build plan",
        f"Over the range we measured, size does not close the gap — and "
        f"{F['BestOpenParams']}B is already past the point where a guard is cheap to run on "
        f"every request, which was the reason to self-host in the first place. Three points, "
        f"one of them non-monotonic, identify no scaling law, so we do not extrapolate a size "
        f"that would close it.",
        color=ACCENT)
d.notes(s, f"""
One claim: scale is not the escape hatch either.

Exact numbers: 4B to 32B is {F['ScaleGain']}, interval {F['ScaleGainCI']}. The remaining gap
to hosted from 32B is {F['GainOverOpen']}, interval {F['GainOverOpenCI']}. So an eightfold
parameter increase bought slightly less than what was still missing.

The 8B row is worth flagging honestly: it came out below 4B, but the interval
({F['EightBvsFourBCI']}) includes zero, so the correct statement is "no gain", not "worse".
Do not oversell it -- but do use it to make the point that size and guard quality are not
tightly coupled.

The build-plan consequence is the part an executive audience should leave with. A 32B guard on
every inbound request is a serving cost close to the thing we were trying to avoid. If we are
paying that, the hosted option deserves another look on cost grounds alone.
""")

# ─────────────────────────────────── 7 · the cascade: evidence for the rule
s = d.blank()
y = d.header(s, "why the rule works", "You do not have to choose. Escalating the uncertain "
             "fifth buys half the difference",
             "Recall against the share of requests sent out, at a fixed 5% false-alarm budget")
picture(s, "exec_cascade", M, y - Inches(0.08), Inches(8.30), Inches(3.40), align="left")
callout(s, M + Inches(8.60), y + Inches(0.00), CW - Inches(8.60), Inches(1.72),
        "the curve is steep early",
        "The first requests escalated are the ones the in-house guard is least sure about, so "
        "they carry the most information. Value per escalated request falls as the share "
        "rises — which is why a fifth buys half.", color=BLUE)
callout(s, M + Inches(8.60), y + Inches(1.90), CW - Inches(8.60), Inches(1.50),
        "it is a dial, not a switch",
        "Legal sets the ceiling on what may leave; budget sets the rest. Pick a point on the "
        "curve — the architecture does not change.", color=GREEN)
callout(s, M, y + Inches(3.52), CW, Inches(0.92), "what this costs",
        f"At a 20% escalated share: about ${float(F['BestCost'])*0.2:.2f} per thousand requests "
        f"instead of ${F['BestCost']}, and the {F['BestMedianMs']} ms hosted round-trip lands "
        f"on one request in five rather than on all of them.", color=AMBER)
d.notes(s, """
This is the evidence slide for the rule on slide 3. One claim: the relationship is smooth and
steep early, so escalation is a dial rather than an architecture decision.

Mechanism, and it is why this works at all: the requests chosen for escalation are the ones
nearest the in-house guard's own decision line -- exactly the cases where a second opinion can
change the answer. Escalating confident rows adds cost and no accuracy.

Method, if challenged. One global false-alarm budget across the whole curve, so every point is
comparable and we are not buying recall by alarming more. The two guards are fused on RANK, not
raw score, because a logit margin and a self-reported 0-100 risk are not on the same scale. We
tried a per-band threshold first and rejected it: with a small escalated slice there are too few
deferred negatives to place a stable threshold, and the curve went non-monotone for purely
numerical reasons. The figure stops at 50% because past that you have outsourced the majority;
the upper rule shows where full escalation lands.

Cost is linear in the escalated share, so a fifth of the traffic is a fifth of the bill.

One rounding note, in case someone has the technical report open beside this. The chart
computes "% of the gap closed" from the raw cascade curve; the report's own gap-ladder figure
computes it from the committed three-decimal values (.856, .787, .896) and therefore prints
+63% at the 30% point where this chart prints 64%. Same measurement, 63.3% against 63.5%, and
the difference is which side of the rounding the inputs sit on. Say that rather than
improvising -- the three annotated points here are internally consistent with each other.
""")

# ────────────────────── 8 · if nothing may leave: combining what we already have
s = d.blank()
y = d.header(s, "if nothing may leave", "Combining our own guards helps a little. It does not "
             "reach the hosted model",
             f"Priced against our strongest single in-house guard ({FN.pct(F['EnsBestSingle'])}) "
             f"— which is itself a five-seed average — with what each costs per request")
# Two corrections live in these rows. (1) The .834 "best single guard" IS the 32B's five
# tuned adapters averaged, so it costs five calls, not one -- and row 2's +0.026 is the
# mechanism that produced it, not a gain to stack on top of it. Reading rows 1 and 2 together
# as "83% + 0.026" gives ~86%, which appears nowhere in the report and would beat the fitted
# 18-guard stack. (2) The cascade row starts from the 3B inline guard at .787, not from this
# row: the measured move is 79% -> 84%, not 83% -> 84%.
rows = [["Option", "Catch rate", "Cost per request", "Verdict"],
        [f"Best single in-house guard — the {F['BestOpenName']}'s 5 tuned adapters, averaged",
         FN.pct(F["EnsBestSingle"]), "5 model calls", ("the baseline below", SLATE)],
        ["Seed-averaging — the mechanism behind that row, not an addition to it",
         f"{F['EnsSeedGain']} over one tuned seed", "no new training",
         ("worth doing", GREEN)],
        [f"Average all {F['EnsMembers']} guards equally", FN.pct(F["EnsCommittee"]),
         f"{F['EnsMembers']} model calls", ("worse than one", ACCENT)],
        [f"Weighted blend of all {F['EnsMembers']}", FN.pct(F["EnsStack"]),
         f"{F['EnsMembers']} calls + labelled data", ("best in-house", BLUE)],
        ["Escalate the uncertain 20% — from the 3B inline guard",
         f"{FN.pct(b(CA['LocalOnly']))} → {FN.pct(CA['Twenty'])}", "1.2 model calls",
         ("nearly matches, far cheaper", BLUE)]]
table(s, M, y + Inches(0.04), CW, rows, [0.35, 0.19, 0.23, 0.23], row_h=Inches(0.54))
callout(s, M, y + Inches(3.32), CW, Inches(1.16), "the finding that matters for planning",
        f"Escalating one request in five gets within about a point of the {F['EnsMembers']}-model "
        f"blend ({FN.pct(CA['Twenty'])} vs {FN.pct(F['EnsStack'])}) at a fifteenth of the compute "
        f"and with no labelled data to collect — measured on top of the 3B inline guard; a "
        f"cascade over a larger inline guard was never measured. Ensembling is the right answer "
        f"only when nothing may leave at all; then it is the in-house ceiling, worth about a "
        f"quarter of the gap.",
        color=AMBER)
d.notes(s, f"""
This slide exists because "just ensemble the small models" is the most common suggestion in the
room, and it deserves a measured answer rather than an opinion.

Two things about the top two rows, because they are easy to misread and an earlier version of
this table invited exactly that. First, the {FN.pct(F['EnsBestSingle'])} baseline is NOT a
single forward pass: the strongest single open configuration we hold is the {F['BestOpenName']}
with its five tuned adapters averaged, so it is five model calls. Second, row 2's
{F['EnsSeedGain']} is measured against a SINGLE TUNED SEED of the same checkpoint, and it is
what produced row 1 -- it is not a gain to add on top of it. Concretely on the
{F['BestOpenName']}: {F['EnsBestSingle']} for the seed ensemble against 0.830 for its own
untuned base, i.e. +0.004 over the base and about +0.026 over one of its own tuned seeds.
Adding {F['EnsSeedGain']} to {FN.pct(F['EnsBestSingle'])} would produce a number that beats the
fitted {F['EnsMembers']}-guard stack and appears nowhere in the report.

"No new training" is the honest cost entry for row 2: the adapters already exist, so it costs
no GPU time to CREATE. It is not free to RUN -- averaging five adapters is five forward passes
per request, in a deck whose whole argument is per-request latency and cost.

Averaging all {F['EnsMembers']} guards equally is WORSE than just using the best one
({F['EnsCommittee']} against {F['EnsBestSingle']}). With members that differ this much in
quality, equal weights dilute the strong ones. If someone proposes "ensemble everything", this
is the row to show them.

A fitted weighted blend does beat any single guard, reaching {F['EnsStack']}, about a quarter of
the way to hosted. Two costs are easy to miss: {F['EnsMembers']} forward passes per request, and
it needs LABELLED in-domain data to fit the weights -- which is usually the thing a team reaching
for an ensemble does not have. We scored it {F['EnsFolds']}-fold out-of-fold so it is not
grading its own homework, and its fitted weights include a negative coefficient on one guard,
which means it is exploiting error structure that may not survive a change of traffic.

The last row is the point: escalating the uncertain fifth reaches {FN.pct(CA['Twenty'])} against the
blend's {FN.pct(F['EnsStack'])} -- it does NOT beat the blend, it very nearly matches it for a
fifteenth of the compute ({F['EnsMembers']} calls against 1.2) and no labelled data. If
challenged on this, the honest line is that the blend is still the in-house ceiling; escalation
is the cheaper route to almost the same place. An earlier version of this slide printed 87%
here and claimed escalation beat the blend, which reversed the ordering. So ensembling is not
the route -- unless lane 1 applies and nothing may leave, in which case it is the in-house
ceiling and worth building.

Be precise about what that last row was measured ON TOP OF, because the table would otherwise
read as "escalation adds one point to our best in-house guard". It does not. It was measured
with the SmolLM3-3B inline guard underneath it, moving {b(CA['LocalOnly'])} to
{b(CA['Twenty'])} -- five points, not one. A cascade sitting on top of a 32B inline guard is a
combination that does not exist anywhere in the report, and an engineer who plans around
"83% to 84%" will build something we never measured.

One in-house repair the table cannot carry, if someone asks whether there is anything else:
averaging a base with its own adapter recovers +0.076 of transfer ranking for one extra pass
and no retraining. That is the cheapest in-house repair in the research -- but it is measured
in macro-AP on the internal panel, not as catch rate at a matched budget on these rows, so it
cannot be dropped into this table without breaking the same-rows, same-budget premise that
makes the table comparable at all. Offer it as a next measurement, not as a fifth row.
""")

# ───────────────────────────── 8b · the lane we cannot escalate, and its evidence
# The deck recommends an architecture built around lane 1 -- borrower-bearing traffic that
# cannot lawfully leave, and therefore cannot be escalated -- and until now said nothing about
# how well the in-house guard actually does on it. The report has a purpose-built instrument
# for exactly that traffic and its findings are unflattering. Saying so does not weaken the
# recommendation (there is no lawful alternative for that traffic); it makes the last slide's
# "build the domain instrument" ask a consequence of a measurement rather than an aspiration.
# Figures are from the frozen v1_hmda2022 mortgage benchmark and its zero-shot baseline table.
s = d.blank()
y = d.header(s, "the lane we cannot escalate",
             "On the one lane where nothing may leave, we have the least evidence our guard "
             "works",
             "Measured on the instrument we built for exactly that traffic: 994 dual-labelled "
             "mortgage rows, scored zero-shot")
cw3 = (CW - Inches(0.6)) / 3
statcard(s, M, y + Inches(0.06), cw3, Inches(1.62), "502 of 994",
         "rows that read safe to a general safety guard and still solicit a compliance "
         "violation. The largest non-benign block, by design — this is the payload.",
         color=ACCENT, value_size=21)
statcard(s, M + cw3 + Inches(0.3), y + Inches(0.06), cw3, Inches(1.62), "0.85  vs  0.555",
         "best zero-shot ranking of policy violations, against the chance floor a coin flip "
         "already scores. The whole band is 0.12–0.30 above chance.",
         color=ACCENT, value_size=21)
statcard(s, M + 2 * (cw3 + Inches(0.3)), y + Inches(0.06), cw3, Inches(1.62),
         "below the median",
         "where all four guards rank the worked violation, against the benign inquiries in "
         "its own split. One of them ranks it below every single one.",
         color=ACCENT, value_size=21)
cw2 = (CW - Inches(0.45)) / 2
callout(s, M, y + Inches(1.86), cw2, Inches(1.66),
        "and we cannot even give you an operating point",
        "The fixed 5%-false-alarm threshold on this benchmark is not reportable: its count of "
        "caught violations swung by more than 50 rows across library versions of one quantile "
        "routine. That is itself a finding about threshold transfer in this domain — but it "
        "means there is no number here to plan a deployment against.", color=ACCENT)
callout(s, M + cw2 + Inches(0.45), y + Inches(1.86), cw2, Inches(1.66),
        "this does not change the recommendation — it prices it",
        "There is no lawful alternative for borrower-bearing traffic, so lane 1 stands. What "
        "changes is what we owe it: the domain instrument on the last slide stops being a "
        "nice-to-have and becomes the direct consequence of a measurement we already have. "
        "Labels here are LLM-judge against written policy cards, not counsel-reviewed.",
        color=GREEN)
d.notes(s, """
Do not skip this slide to save time, and do not let it be discovered in Q&A instead. The whole
architecture is built around lane 1, lane 1 is the lane where escalation is unavailable by
construction, and this is the only evidence we have about how the in-house guard behaves there.

The four facts, in the order they land hardest. The G0/D1 stratum -- reads safe to a general
guard, is a mortgage-policy violation -- is 502 of 994 rows, which is the point of the
benchmark rather than an accident of sampling. Threshold-free ranking of those violations tops
out at 0.85 against a 0.555 chance floor set by the split's own base rate (81 of 146 rows are
policy-positive), so the observed band is 0.12 to 0.30 above chance, not 0.85 worth of skill.
On the worked row, all four zero-shot guards rank a coded redlining-by-proxy request below the
median benign mortgage inquiry in the same split, and one of them ranks it below all 65. And
the fixed operating point is deliberately not tabulated, because its catch count moved by more
than 50 rows across library versions of a quantile routine -- knife-edge on clustered scores.

The tone to take: this is the instrument working. We built a measuring stick for the traffic we
care most about, held it to the same standard as the guards, and it says the guards miss the
payload. That is a reason to fund the instrument, not a reason to distrust the recommendation --
there is no lawful alternative for prompts carrying borrower data, so the choice on that lane
was never between a good option and a better one.

Scope it honestly if a compliance function is in the room: these labels are LLM-judge against
written policy cards with self-consistency checks, no subject-matter-expert adjudication, no
human agreement statistic, and the public-test split is a single 146-row set. It surfaces guard
behaviour. It certifies nothing, and it licenses no fair-lending conclusion.
""")

# ──────────────────────────────────────────────────── 9 · what it costs to host
s = d.blank()
# The refusal finding is the strongest non-obvious point in the deck -- a compliance control
# that intermittently declines to answer, non-reproducibly, for reasons the operator cannot
# inspect or appeal -- and it was a footnote on a costs slide. It is now the headline, and the
# invoice costs are the supporting table rather than the other way round.
y = d.header(s, "what hosted costs",
             "The bill and the latency are the easy part. The control sometimes declines to "
             "answer",
             "Four costs, three of which do not appear on an invoice")
cw = (CW - Inches(0.45)) / 2
rows = [["", "Self-hosted small", "Hosted frontier"],
        ["Median latency", "10–25 ms", f"~{F['BestMedianMs']} ms"],
        # "GPU you own" is the honest entry and it is not a price. We have not amortised our
        # own side; the notes carry that answer rather than leaving it to be improvised.
        ["Cost per 1k prompts", "GPU you own — not amortised", f"${F['BestCost']} list"],
        ["Prompt leaves our network", "no", "yes"],
        ["We choose the alarm rate", "exactly", "approximately"]]
table(s, M, y + Inches(0.10), cw, rows, [0.44, 0.28, 0.28])
callout(s, M + cw + Inches(0.45), y + Inches(0.10), cw, Inches(1.62),
        "the one we did not expect",
        "The provider refused to evaluate a subset of prompts outright, and did so "
        "inconsistently — the same prompt was refused on some runs and not others. A "
        "guardrail that intermittently declines to answer is itself an audit finding.",
        color=ACCENT)
callout(s, M + cw + Inches(0.45), y + Inches(1.86), cw, Inches(1.50),
        "control over the alarm rate",
        "Our own model exposes a continuous score, so we can set the false-alarm budget "
        "precisely. The hosted model reports a coarse 0–100 value, so we can only land "
        "near a target.", color=BLUE)
callout(s, M, y + Inches(3.46), CW, Inches(1.06), "not a reason to refuse hosted",
        "It is a reason to scope it: hosted on traffic where the prompt may leave and "
        "latency is not critical; self-hosted on the regulated, latency-sensitive path.",
        color=GREEN)
d.notes(s, f"""
Work down the table, then land on the two callouts, which are the non-obvious costs.

Latency: {F['BestMedianMs']} ms median measured under load at high concurrency, so treat it as
an upper bound for a single request; the small-model figure is a batched GPU measurement. Even
allowing for that, the ratio is roughly {F['Slowdown']}x.

Cost: billed tokens at public list prices. It excludes the amortised GPU on the self-hosted
side, which is exactly why the comparison is about where to spend rather than a like-for-like
price.

Somebody will ask "so what does OURS cost per thousand?" -- have the honest answer ready rather
than improvising one. We have not amortised it. What we can quote is the physical measurement
behind it: batched per-row inference at batch 16 on a single A100, 10-25 ms per call at P50, so
a thousand prompts is on the order of tens of seconds of one GPU. Turning that into a unit
price needs a utilisation assumption and a capital-recovery period that finance owns, not us.
Offer to come back with a number rather than inventing one at the table.

The refusal finding is the one to say slowly, and it is now the headline of this slide. We saw a set of prompts refused by the provider's
own input filter before the model ever saw them, and the refusals were not reproducible --
different runs, different outcomes, no row refused every time. For a control that has to be
auditable, "sometimes declines, for reasons we cannot inspect or appeal" is a compliance
property, not a technical footnote.

Close on the scoping sentence. The answer is a split, and this slide is why.
""")

# ────────────────────────────────────────────── 10 · the useful surprise
s = d.blank()
y = d.header(s, "the useful finding", "If we must self-host, a bigger untuned model beats a "
             "tuned small one",
             "Because tuning's cost lands on exactly the traffic a guardrail exists to catch")
cw = (CW - Inches(0.45)) / 2
bullets(s, M, y + Inches(0.06), cw, Inches(2.9), [
    ("Tuning a small model. ", "Large gain on familiar traffic, and it gives back accuracy "
     "on unfamiliar traffic — which is where novel attacks arrive."),
    ("A bigger untuned model. ", "Recovers most of that same gain on familiar traffic, and "
     "holds its accuracy on unfamiliar traffic."),
    ("So for a guardrail specifically. ", "Prefer spending on a stronger base model over "
     "spending on a tuning programme."),
], size=14.5)
callout(s, M + cw + Inches(0.45), y + Inches(0.06), cw, Inches(1.70),
        "why this is worth money",
        "It reverses the intuitive plan. Tuning looks like the cheap, targeted option; "
        "measured on held-out traffic it is the one that quietly costs you coverage.",
        color=GREEN)
callout(s, M + cw + Inches(0.45), y + Inches(1.92), cw, Inches(1.70),
        "confidence level",
        "This is a comparison between different model sizes, not a controlled experiment at "
        "fixed size. Treat it as a strong steer on where to spend, not a proven law.",
        color=AMBER)
d.notes(s, """
This slide did not exist in the plan. It fell out of the scale experiment and is probably the
most commercially useful thing in the deck.

The mechanism: fine-tuning specialises. It buys accuracy on the sources you trained on and
withdraws it from sources you did not. A bigger base model arrives with broad competence
already and does not have to trade any of it away. For a guardrail, the traffic you did not
anticipate is precisely the traffic that matters, so a technique that trades away
unfamiliar-traffic accuracy is badly matched to the job.

Numbers are in the technical report; the honest caveat is on the slide. Comparing a 32B untuned
model to a tuned 4B is a deployment-choice comparison, not a controlled one -- the 32B costs
much more to serve. What we can say is that the tuning route did not close the vendor gap and
made held-out accuracy worse on most models, while the size route at least moved in the right
direction on both regimes.
""")

# ─────────────────────────────────────────────────── 11 · what we are unsure of
s = d.blank()
y = d.header(s, "confidence and limits", "What this does not tell you",
             "Stated plainly, so the recommendation can be trusted where it does apply")
bullets(s, M, y + Inches(0.06), CW, Inches(3.2), [
    ("It is one benchmark. ", "Expert-annotated finance, healthcare and law prompt safety. "
     "It is the strongest external evidence we have, and it is still a single instrument — "
     "and our own research shows guard rankings reorder when the benchmark changes."),
    ("It is prompt-only. ", "We judged incoming prompts, not the assistant's replies, and not "
     "the harder mortgage-compliance question of whether a specific answer breaks a specific "
     "rule — which, as slide 10 shows, is the lane where our own instrument says the guards "
     "miss the payload."),
    ("Vendor behaviour is a moving target. ", "Model versions change under you. The "
     "self-hosted side is pinned to an exact revision and reproduces indefinitely; the hosted "
     "side does not."),
    ("Prices are list, not invoices. ", "Cost figures are billed tokens at public rates and "
     "exclude the GPU cost on our own side."),
    # The recommendation on slide 3 rests on the escalation curve, and the report is explicit
    # that the curve is optimistically tuned. Stating it only in the notes let the deck carry
    # the number at a higher confidence than the paper allows.
    ("The escalation curve is optimistically tuned. ", "Its decision line and its global 5% "
     "budget are both chosen on the same rows it is then scored on, so a live deployment "
     "should expect to close somewhat less than 51% of the gap at a fifth escalated."),
    # The research report re-reads its own headline inside the alarm budget; an exec deck that
    # quotes catch rates at a 5% budget should say that the underlying trade was measured that
    # way too, otherwise it inherits an average-ranking number without saying so.
    ("Averages hide the operating point. ", "Every catch rate here is already at a 5% "
     "false-alarm budget, but the underlying fine-tuning research reports an average over the "
     "whole score ranking. Re-scored inside the budget the direction is unchanged and the "
     f"effect is {FN.bare(LFP['LowFprTransAmplification'])}x larger — so treat average-metric "
     "results from elsewhere as understating what you would see in production."),
], size=12.2, gap=8)
callout(s, M, y + Inches(3.35), CW, Inches(1.08), "what would change the recommendation",
        "A regulated-domain benchmark that scores whole answers against specific rules, with "
        "expert adjudication. That is the instrument we do not have, and building it is the "
        "highest-value next step.", color=BLUE)
d.notes(s, """
Do not rush this slide, and do not apologise through it. Naming the limits is what makes the
recommendation on slide 3 credible.

The first bullet is the most important and the most self-undermining, which is why it belongs
here: the underlying research finding is that guard rankings depend on the benchmark. We are
handing you a recommendation derived from one benchmark. It is external and expert-annotated,
which is the best tier available to us, but a different instrument could reorder the middle of
the table. It is unlikely to reverse a ten-point gap.

The pinning point matters more than it sounds. A hosted model can change under a fixed
deployment; that is a change-control problem for a compliance control. Our own checkpoints are
pinned to a content hash and will reproduce in a year.

Close on the final callout: the thing worth funding is the measuring instrument for our actual
domain, not more guard comparisons on general safety.
""")

# ────────────────────────────────────────────────────────── 12 · what we would do
# Closing slide shares the title slide's darker surface, bookending the deck.
s = d.blank(title_slide=True)
y = d.header(s, "proposal", "What we would do next",
             "Sequenced so each step is cheap and the expensive one is last")
cw = (CW - Inches(0.45)) / 2
bullets(s, M, y + Inches(0.06), cw, Inches(3.3), [
    ("Now · keep the in-house guard inline on everything. ", "It already runs, it answers in "
     "~20 ms, and it is the only lane available for regulated traffic. No new build."),
    ("Now · add the escalation lane. ", "Route the least-confident band to the hosted model "
     "where the data may lawfully leave. Start at 10–20% and read the curve."),
    ("Now · treat a provider refusal as 'needs review'. ", "Never as 'safe'. A one-line "
     "policy change, and today it is a silent gap."),
    ("Next · stop the tuning workstream as a gap-closer. ", "Keep it only for rescuing weak "
     "models. Redirect the effort."),
    ("Then · build the domain instrument. ", "A mortgage-compliance benchmark that scores "
     "whole answers against specific rules, with expert adjudication. Slide 10 is why: on the "
     "lane we cannot escalate, the instrument we have says the guards miss the payload."),
], size=13.5, gap=10)
callout(s, M + cw + Inches(0.45), y + Inches(0.06), cw, Inches(1.20),
        "the cheapest item is the refusal one",
        "It costs a policy change, and it closes a case where an unsafe prompt currently "
        "receives no verdict at all.", color=GREEN)
callout(s, M + cw + Inches(0.45), y + Inches(1.36), cw, Inches(1.32),
        "decision one — a legal call, not a technical one",
        "Whether borrower-bearing prompts may leave our network at all. That single answer "
        "decides how much of our traffic can use the more accurate option.", color=ACCENT)
# Slide 8 correctly frames the escalated share as a standing tradeoff between the legal
# ceiling and the budget rather than a one-time architecture choice. A standing tradeoff needs
# a standing owner, and naming one is the second thing only this room can settle.
callout(s, M + cw + Inches(0.45), y + Inches(2.84), cw, Inches(1.32),
        "decision two — who owns the dial",
        "The escalated share is not set once. Someone has to own moving it as the legal "
        "ceiling and the budget move, and to review it on a schedule.", color=ACCENT)
d.notes(s, """
Two minutes. The first two items need no budget. The third frees capacity. The fourth is the
only one that asks for money.

Push for BOTH decisions before the meeting ends. The first is whether prompts containing
borrower information may go to a third-party API; everything about the split architecture
depends on that answer, and it is not ours to make. The second is who owns the escalation dial
once it exists -- slide 8 shows the escalated share is a smooth tradeoff between the legal
ceiling and the budget, which means it is a standing decision with a standing owner rather than
an architecture we choose once and forget. If nobody owns it, it will sit at whatever we set on
day one until an incident moves it.

The refusal item is the one to insist on regardless of the rest. Today, when the provider
declines to evaluate a prompt, our pipeline records no verdict. That is indistinguishable from
a pass unless we say otherwise, and the prompts most likely to be refused are the ones most
likely to be genuinely unsafe.

If asked for the whole thing in one sentence: hosted is about ten points better at catching
unsafe prompts and we should use it wherever the data is allowed to travel, we should stop
expecting to close that gap ourselves, and the thing worth building is a benchmark for our own
domain rather than another guard.
""")

path = d.save()
print(f"wrote {path}  ({d.n + 1} slides)")
